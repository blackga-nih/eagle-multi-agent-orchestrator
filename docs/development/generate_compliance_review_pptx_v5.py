"""Generate the v5 PowerPoint deck for the compliance officer's template review.

v5 differences from v4:
- New slide right after the title: "Templates to Verify" — single-page list of all
  eight templates with version + KB-grounded badge so the reviewer can see the
  full scope before reading per-template detail.
- COR designation slide now surfaces the actual KB source filename
  (NIH COR Appointment Memorandum.docx). v4 said only "NIH COR Handbook + FAQ" —
  the user pointed out we should call out that the KB file is the
  Appointment Memorandum (the NCI/NIH name for the same FAR 1.602-2(d)
  designation instrument).
- "What changed since v1" expanded with a v5 entry documenting the
  Appointment Memorandum reconciliation.
"""

from pathlib import Path
import re
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

REPO_ROOT = Path(r"C:\Users\blackga\Desktop\eagle\sm_eagle")
TEMPLATES_DIR = REPO_ROOT / "docs/development/templates"
OUT_PPTX = REPO_ROOT / "docs/development/20260505-171152-report-compliance-template-review-v5.pptx"

NCI_NAVY = RGBColor(0x00, 0x33, 0x66)
NCI_LIGHT = RGBColor(0xF6, 0xF8, 0xFB)
TEXT_DARK = RGBColor(0x1C, 0x1C, 0x1C)
TEXT_GREY = RGBColor(0x55, 0x55, 0x55)
ACCENT_LINE = RGBColor(0xB0, 0xC0, 0xD0)
KB_GREEN = RGBColor(0x16, 0x6B, 0x3A)
KB_AMBER = RGBColor(0xB3, 0x6B, 0x00)


# slug, version, doc_name, kb_grounded, kb_label
TEMPLATES = [
    ("cor_designation", "v2", "COR Designation Memo / Appointment Memorandum", True,
     "NIH COR Appointment Memorandum.docx + NIH COR Handbook + FAQ"),
    ("sb_review", "v2", "HHS-653 Small Business Review", True, "HHS AA 2023-02 Am. 4"),
    ("section_889", "v1", "Section 889 Compliance", False, "(no KB source - FAR 4.21 only)"),
    ("section_508", "v2", "Section 508 Compliance", True, "OAG FY25-02 + CD 2024-01"),
    ("priority_sources_checklist", "v2", "Required & Priority Sources Checklist", True, "NIH Reference + HHSAM 308"),
    ("subk_review", "v2", "Subcontracting Plan Review", True, "HHS SubK Review Process SOP"),
    ("qasp", "v1", "Quality Assurance Surveillance Plan", False, "(no KB source - FAR 46.401 only)"),
    ("source_selection_plan", "v2", "Source Selection Plan (SSP)", True, "HHS Down Select Guides"),
]


def parse_template(slug: str, version: str) -> dict:
    """Pull purpose / authority / audience / sections / worked example from a template MD."""
    fp = TEMPLATES_DIR / f"{slug}-template-{version}.md"
    text = fp.read_text(encoding="utf-8")

    def grab(label: str) -> str:
        m = re.search(rf"\*\*{label}\*\*[:\s]+(.*?)(?=\n\n|\n\*\*|\n## )", text, re.DOTALL)
        return (m.group(1).strip() if m else "").replace("\n", " ")[:400]

    sec_match = re.search(r"## Required Sections\s*\n(.*?)(?=\n## )", text, re.DOTALL)
    sections: list[tuple[str, str]] = []
    if sec_match:
        for line in sec_match.group(1).splitlines():
            line = line.strip()
            m = re.match(r"^\d+\.\s+\*?\*?(.*?)\*?\*?\s*[—-]\s+(.*)$", line)
            if m:
                sections.append((m.group(1).strip(), m.group(2).strip()))
            else:
                m = re.match(r"^\d+\.\s+(.*)$", line)
                if m:
                    parts = re.split(r"\s[—-]\s", m.group(1), maxsplit=1)
                    if len(parts) == 2:
                        sections.append((parts[0].strip(" *"), parts[1].strip()))
                    else:
                        sections.append((m.group(1)[:60], ""))

    we_match = re.search(r"##\s+(?:Small\s+)?Worked Example.*?\n(.*?)(?=\n## )", text, re.DOTALL)
    worked = ""
    if we_match:
        we = we_match.group(1).strip()
        worked = re.sub(r"\n\|.*", "", we, flags=re.DOTALL).strip()
        worked = worked.replace("**", "")[:850]

    return {
        "purpose": grab("Purpose"),
        "authority": grab("Authority"),
        "audience": grab("Audience"),
        "sections": sections[:11],
        "worked": worked,
    }


def navy_band(slide, prs) -> None:
    band = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), prs.slide_width, Inches(0.4))
    band.fill.solid()
    band.fill.fore_color.rgb = NCI_NAVY
    band.line.fill.background()


def slide_title(slide, title_text: str) -> None:
    tb = slide.shapes.add_textbox(Inches(0.7), Inches(0.55), Inches(12), Inches(0.7))
    p = tb.text_frame.paragraphs[0]
    r = p.add_run()
    r.text = title_text
    r.font.size = Pt(24)
    r.font.bold = True
    r.font.color.rgb = NCI_NAVY


prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)


# 1. Title slide
s = prs.slides.add_slide(prs.slide_layouts[6])
navy_band(s, prs)
tb = s.shapes.add_textbox(Inches(0.7), Inches(2.4), Inches(12), Inches(1.2))
r = tb.text_frame.paragraphs[0].add_run()
r.text = "Compliance Template Review - v2"
r.font.size = Pt(40); r.font.bold = True; r.font.color.rgb = NCI_NAVY
sub = s.shapes.add_textbox(Inches(0.7), Inches(3.5), Inches(12), Inches(0.8))
sr = sub.text_frame.paragraphs[0].add_run()
sr.text = "Eight Acquisition Document Templates - KB-Grounded Drafts"
sr.font.size = Pt(20); sr.font.color.rgb = TEXT_DARK

meta = s.shapes.add_textbox(Inches(0.7), Inches(5.4), Inches(12), Inches(1.5))
mtf = meta.text_frame
for i, line in enumerate([
    "EAGLE Platform Team - NCI Office of Acquisitions",
    "Date: 2026-05-05",
    "Audience: Reviewing Compliance Officer",
    "v5 deck: scope summary up front; COR Appointment Memorandum reconciled with cor_designation template",
]):
    p = mtf.paragraphs[0] if i == 0 else mtf.add_paragraph()
    rr = p.add_run()
    rr.text = line
    rr.font.size = Pt(14)
    rr.font.color.rgb = TEXT_GREY
    if i == 3:
        rr.font.bold = True
        rr.font.color.rgb = KB_GREEN


# 2. NEW: Templates to Verify (scope summary up front)
s = prs.slides.add_slide(prs.slide_layouts[6])
navy_band(s, prs)
slide_title(s, "Templates to Verify (8)")

intro = s.shapes.add_textbox(Inches(0.7), Inches(1.25), Inches(12), Inches(0.55))
itf = intro.text_frame; itf.word_wrap = True
ip = itf.paragraphs[0]
ir = ip.add_run()
ir.text = ("This deck asks the compliance officer to verify the following eight EAGLE prompt "
          "templates. KB column shows whether the v2 draft was grounded in the EAGLE knowledge base.")
ir.font.size = Pt(12); ir.font.color.rgb = TEXT_GREY

headers = ["#", "Template", "Slug", "Ver", "KB", "Source"]
col_w = [Inches(0.4), Inches(3.6), Inches(2.3), Inches(0.6), Inches(0.6), Inches(4.8)]
table = s.shapes.add_table(len(TEMPLATES) + 1, len(headers),
                           Inches(0.5), Inches(1.95),
                           sum((c for c in col_w), Inches(0)), Inches(4.6)).table
for ci, w in enumerate(col_w):
    table.columns[ci].width = w

for ci, htext in enumerate(headers):
    c = table.cell(0, ci); c.text = ""
    p = c.text_frame.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    r = p.add_run(); r.text = htext
    r.font.size = Pt(11); r.font.bold = True; r.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
    c.fill.solid(); c.fill.fore_color.rgb = NCI_NAVY

for ri, (slug, version, doc_name, kb_grounded, kb_label) in enumerate(TEMPLATES, 1):
    kb_mark = "Yes" if kb_grounded else "No"
    kb_color = KB_GREEN if kb_grounded else KB_AMBER
    cells = [str(ri), doc_name, slug, version, kb_mark, kb_label]
    for ci, val in enumerate(cells):
        c = table.cell(ri, ci); c.text = ""
        p = c.text_frame.paragraphs[0]
        if ci in (0, 3, 4):
            p.alignment = PP_ALIGN.CENTER
        r = p.add_run(); r.text = val
        r.font.size = Pt(10); r.font.color.rgb = TEXT_DARK
        if ci == 1:
            r.font.bold = True; r.font.color.rgb = NCI_NAVY
        if ci == 4:
            r.font.bold = True; r.font.color.rgb = kb_color
        if ri % 2 == 0:
            c.fill.solid(); c.fill.fore_color.rgb = NCI_LIGHT

footer = s.shapes.add_textbox(Inches(0.5), Inches(6.7), Inches(12.5), Inches(0.5))
fp = footer.text_frame.paragraphs[0]
fr = fp.add_run()
fr.text = ("Per-template detail starts after the change summary. Approval matrix is the second-to-last slide; "
          "use it to record APPROVED / APPROVED W/ CHANGES / REJECTED for each row above.")
fr.font.size = Pt(10); fr.font.italic = True; fr.font.color.rgb = TEXT_GREY


# 3. What changed since v1
s = prs.slides.add_slide(prs.slide_layouts[6])
navy_band(s, prs)
slide_title(s, "What changed since v1")

body = s.shapes.add_textbox(Inches(0.7), Inches(1.3), Inches(12), Inches(5.8))
tf = body.text_frame; tf.word_wrap = True

bullets = [
    "v1 was drafted from generic FAR/HHSAR knowledge - WITHOUT consulting the EAGLE knowledge base.",
    "Spot check identified the gap; SSO refresh enabled S3 enumeration of the approved/ prefix.",
    "KB sources found for 6 of 8 templates; redrafted as v2 with explicit grounding.",
    "Two templates (qasp, section_889) had no KB source - kept at v1; flagged for officer.",
    "Each v2 has a Source Grounding table at the bottom mapping content to its KB doc.",
    ("v5 deck: cor_designation reconciled with the KB file 'NIH COR Appointment Memorandum.docx' "
     "in supervisor-core/essential-templates/ - same FAR 1.602-2(d) instrument, NCI/NIH name."),
]
for i, text in enumerate(bullets):
    p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
    p.space_after = Pt(10)
    rb = p.add_run(); rb.text = "*  " + text
    rb.font.size = Pt(15); rb.font.color.rgb = TEXT_DARK
    if i == len(bullets) - 1:
        rb.font.bold = True
        rb.font.color.rgb = KB_GREEN


# 4. Material differences caught
s = prs.slides.add_slide(prs.slide_layouts[6])
navy_band(s, prs)
slide_title(s, "Material differences v1 missed")

diffs = [
    ("cor_designation", "Six prescribed responsibility categories from NIH COR Handbook Appendix 8A; FAC-COR dollar tiers (I/II/III); separate signature page; reconciled to KB file 'NIH COR Appointment Memorandum.docx'"),
    ("sb_review", "HHS AA 2023-02 Amendment 4 5-band threshold matrix; SBCX submission; 7-day SBS / 12-day PCR review windows"),
    ("section_508", "Replaced superseded HHSAR 352.239-73/74 with operative CD 2024-01 clauses; corrected WCAG baseline 2.0 -> 2.1 AA per OAG FY25-02"),
    ("priority_sources_checklist", "HHSAM 308.104 three-tier framework (Required Use -> Mandatory Use -> Mandatory Consideration); HCA/SPE exception routing"),
    ("subk_review", "Threshold $750K -> $900K; FAR 10-element checklist -> HHS SOP 15-element; three-step CO -> OSDBU -> PCR signature flow"),
    ("source_selection_plan", "HHS Down Select confidence-based rating (High/Some/Low); HHS team roles; RemedyBiz/AT&T case-law on SSDD documentation"),
]

table = s.shapes.add_table(len(diffs)+1, 2, Inches(0.5), Inches(1.3), Inches(12.3), Inches(5.5)).table
table.columns[0].width = Inches(3.2); table.columns[1].width = Inches(9.1)
for ci, htext in enumerate(["Slug", "What v2 added"]):
    c = table.cell(0, ci); c.text = ""
    p = c.text_frame.paragraphs[0]; rr = p.add_run(); rr.text = htext
    rr.font.size = Pt(11); rr.font.bold = True; rr.font.color.rgb = RGBColor(0xFF,0xFF,0xFF)
    c.fill.solid(); c.fill.fore_color.rgb = NCI_NAVY

for ri, (slug, diff) in enumerate(diffs, 1):
    c0 = table.cell(ri, 0); c0.text = ""
    p0 = c0.text_frame.paragraphs[0]; r0 = p0.add_run(); r0.text = slug
    r0.font.size = Pt(10); r0.font.bold = True; r0.font.color.rgb = NCI_NAVY
    if ri % 2 == 0:
        c0.fill.solid(); c0.fill.fore_color.rgb = NCI_LIGHT
    c1 = table.cell(ri, 1); c1.text = ""
    p1 = c1.text_frame.paragraphs[0]; r1 = p1.add_run(); r1.text = diff
    r1.font.size = Pt(10); r1.font.color.rgb = TEXT_DARK
    if ri % 2 == 0:
        c1.fill.solid(); c1.fill.fore_color.rgb = NCI_LIGHT


# 5. Section header
s = prs.slides.add_slide(prs.slide_layouts[6])
navy_band(s, prs)
tb = s.shapes.add_textbox(Inches(0.7), Inches(2.5), Inches(12), Inches(2))
r = tb.text_frame.paragraphs[0].add_run()
r.text = "Eight Templates"
r.font.size = Pt(36); r.font.bold = True; r.font.color.rgb = NCI_NAVY
sp = tb.text_frame.add_paragraph(); sp.space_before = Pt(12)
sr2 = sp.add_run()
sr2.text = "One slide per template - review structure, sections, and KB grounding"
sr2.font.size = Pt(18); sr2.font.color.rgb = TEXT_GREY


# 6-13. Per-template slides
for idx, (slug, version, doc_name, kb_grounded, kb_label) in enumerate(TEMPLATES, 1):
    parsed = parse_template(slug, version)
    s = prs.slides.add_slide(prs.slide_layouts[6])
    navy_band(s, prs)
    slide_title(s, f"Template {idx} of 8 - {doc_name}")

    badge = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(11.0), Inches(0.55), Inches(2.0), Inches(0.4))
    badge.fill.solid()
    badge.fill.fore_color.rgb = KB_GREEN if kb_grounded else KB_AMBER
    badge.line.fill.background()
    btf = badge.text_frame; btf.margin_top = Pt(2); btf.margin_bottom = Pt(2)
    bp = btf.paragraphs[0]; bp.alignment = PP_ALIGN.CENTER
    br = bp.add_run()
    br.text = f"{version.upper()} - KB" if kb_grounded else f"{version.upper()} - model only"
    br.font.size = Pt(10); br.font.bold = True; br.font.color.rgb = RGBColor(0xFF,0xFF,0xFF)

    meta_box = s.shapes.add_textbox(Inches(0.7), Inches(1.25), Inches(12), Inches(1.1))
    mt = meta_box.text_frame; mt.word_wrap = True
    rows = [
        ("slug", slug),
        ("authority", parsed["authority"][:200] or "(see template body)"),
        ("audience", parsed["audience"][:200] or "(see template body)"),
        ("KB source", kb_label),
    ]
    for i, (k, v) in enumerate(rows):
        p = mt.paragraphs[0] if i == 0 else mt.add_paragraph()
        p.space_after = Pt(1)
        rk = p.add_run(); rk.text = f"{k}:  "
        rk.font.size = Pt(10); rk.font.bold = True; rk.font.color.rgb = NCI_NAVY
        rv = p.add_run(); rv.text = v
        rv.font.size = Pt(10); rv.font.color.rgb = TEXT_DARK

    sec_box = s.shapes.add_textbox(Inches(0.7), Inches(2.55), Inches(12.0), Inches(4.4))
    st = sec_box.text_frame; st.word_wrap = True
    h = st.paragraphs[0]; hr = h.add_run()
    hr.text = "Required Sections"
    hr.font.size = Pt(14); hr.font.bold = True; hr.font.color.rgb = NCI_NAVY
    for n, (sec_name, sec_desc) in enumerate(parsed["sections"], 1):
        p = st.add_paragraph(); p.space_before = Pt(3)
        rn = p.add_run(); rn.text = f"{n}. "
        rn.font.size = Pt(10); rn.font.bold = True; rn.font.color.rgb = NCI_NAVY
        rname = p.add_run(); rname.text = (sec_name[:80] + " — ") if sec_name else ""
        rname.font.size = Pt(10); rname.font.bold = True; rname.font.color.rgb = TEXT_DARK
        rdesc = p.add_run(); rdesc.text = sec_desc[:280]
        rdesc.font.size = Pt(10); rdesc.font.color.rgb = TEXT_GREY

    footer = s.shapes.add_textbox(Inches(0.7), Inches(7.05), Inches(12.0), Inches(0.35))
    ft = footer.text_frame
    fp = ft.paragraphs[0]
    fr = fp.add_run()
    fr.text = f"Full template (with worked example, rules, and source grounding):  docs/development/templates/{slug}-template-{version}.md"
    fr.font.size = Pt(10); fr.font.italic = True; fr.font.color.rgb = TEXT_GREY


# 14. Approval matrix
s = prs.slides.add_slide(prs.slide_layouts[6])
navy_band(s, prs)
slide_title(s, "Approval Matrix")

headers = ["#", "Template", "Ver", "APPROVED", "W/ CHANGES", "REJECTED", "Comments"]
col_w = [Inches(0.4), Inches(4.0), Inches(0.7), Inches(1.3), Inches(1.3), Inches(1.3), Inches(2.6)]
table = s.shapes.add_table(len(TEMPLATES)+1, len(headers), Inches(0.5), Inches(1.3),
                            sum((c for c in col_w), Inches(0)), Inches(5.4)).table
for ci, w in enumerate(col_w):
    table.columns[ci].width = w
for ci, htext in enumerate(headers):
    c = table.cell(0, ci); c.text = ""
    p = c.text_frame.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    r = p.add_run(); r.text = htext
    r.font.size = Pt(10); r.font.bold = True; r.font.color.rgb = RGBColor(0xFF,0xFF,0xFF)
    c.fill.solid(); c.fill.fore_color.rgb = NCI_NAVY

for ri, (slug, version, doc_name, kb_grounded, _kb) in enumerate(TEMPLATES, 1):
    ver_label = f"{version} KB" if kb_grounded else f"{version} (!)"
    for ci, val in enumerate([str(ri), doc_name, ver_label, "", "", "", ""]):
        c = table.cell(ri, ci); c.text = ""
        p = c.text_frame.paragraphs[0]
        if ci != 1: p.alignment = PP_ALIGN.CENTER
        r = p.add_run(); r.text = val
        r.font.size = Pt(9); r.font.color.rgb = TEXT_DARK
        if ri % 2 == 0:
            c.fill.solid(); c.fill.fore_color.rgb = NCI_LIGHT

sig = s.shapes.add_textbox(Inches(0.5), Inches(6.85), Inches(12), Inches(0.5))
sp = sig.text_frame.paragraphs[0]
sr = sp.add_run()
sr.text = "Reviewer signature: ___________________________   Name / Title: ___________________________   Date: __________"
sr.font.size = Pt(11); sr.font.color.rgb = TEXT_GREY


# 15. What happens after approval
s = prs.slides.add_slide(prs.slide_layouts[6])
navy_band(s, prs)
slide_title(s, "What happens after approval")
body = s.shapes.add_textbox(Inches(0.7), Inches(1.3), Inches(12), Inches(5.8))
tf = body.text_frame; tf.word_wrap = True
steps = [
    "Approved templates -> committed to server/app/doc_prompts.py as named constants.",
    "Each slug's DocSpec.system_prompt updated to reference the new constant in server/app/doc_registry.py.",
    "Module-load validator confirms registry consistency; unit tests run.",
    "If green, changes ship in the next deployment.",
    "APPROVED W/ CHANGES return to the platform team with the change list, then back here for re-review.",
    "REJECTED - REWORK gets a fresh draft and re-enters this review cycle.",
    "For qasp and section_889: if you identify an authoritative NCI/HHS source we missed, point us to it and we'll redraft with KB grounding.",
]
for i, step in enumerate(steps):
    p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
    p.space_after = Pt(8)
    rb = p.add_run(); rb.text = f"{i+1}.  " + step
    rb.font.size = Pt(14); rb.font.color.rgb = TEXT_DARK


prs.save(str(OUT_PPTX))
print(f"wrote {OUT_PPTX}")
print(f"size: {OUT_PPTX.stat().st_size:,} bytes")
print(f"slides: {len(prs.slides)}")
