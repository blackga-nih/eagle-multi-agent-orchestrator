"""Generate the PowerPoint deck for the compliance officer's template review.

Run from anywhere — paths are absolute. Outputs the .pptx alongside the
markdown report in docs/development/.
"""

from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

REPO_ROOT = Path(r"C:\Users\blackga\Desktop\eagle\sm_eagle")
OUT_PPTX = REPO_ROOT / "docs/development/20260505-171152-report-compliance-template-review-v1.pptx"

NCI_NAVY = RGBColor(0x00, 0x33, 0x66)
NCI_LIGHT = RGBColor(0xF6, 0xF8, 0xFB)
TEXT_DARK = RGBColor(0x1C, 0x1C, 0x1C)
TEXT_GREY = RGBColor(0x55, 0x55, 0x55)
ACCENT_LINE = RGBColor(0xB0, 0xC0, 0xD0)


def add_title_slide(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    _add_navy_header_band(slide)

    # Title
    tb = slide.shapes.add_textbox(Inches(0.7), Inches(2.4), Inches(12), Inches(1.5))
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    r = p.add_run()
    r.text = "Compliance Template Review"
    r.font.size = Pt(40)
    r.font.bold = True
    r.font.color.rgb = NCI_NAVY

    sub = slide.shapes.add_textbox(Inches(0.7), Inches(3.6), Inches(12), Inches(0.8))
    sub_p = sub.text_frame.paragraphs[0]
    sub_r = sub_p.add_run()
    sub_r.text = "Eight Acquisition Document Templates — Approval Requested"
    sub_r.font.size = Pt(20)
    sub_r.font.color.rgb = TEXT_DARK

    meta = slide.shapes.add_textbox(Inches(0.7), Inches(5.6), Inches(12), Inches(1.5))
    mtf = meta.text_frame
    for line in [
        "EAGLE Platform Team — NCI Office of Acquisitions",
        "Date: 2026-05-05",
        "Audience: Reviewing Compliance Officer",
    ]:
        para = mtf.add_paragraph() if mtf.paragraphs[0].text else mtf.paragraphs[0]
        rr = para.add_run()
        rr.text = line
        rr.font.size = Pt(14)
        rr.font.color.rgb = TEXT_GREY


def _add_navy_header_band(slide) -> None:
    from pptx.enum.shapes import MSO_SHAPE
    band = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), prs.slide_width, Inches(0.45))
    band.fill.solid()
    band.fill.fore_color.rgb = NCI_NAVY
    band.line.fill.background()


def add_section_header_slide(prs: Presentation, title: str, subtitle: str = "") -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_navy_header_band(slide)
    tb = slide.shapes.add_textbox(Inches(0.7), Inches(2.5), Inches(12), Inches(2))
    tf = tb.text_frame
    p = tf.paragraphs[0]
    r = p.add_run()
    r.text = title
    r.font.size = Pt(36)
    r.font.bold = True
    r.font.color.rgb = NCI_NAVY
    if subtitle:
        sp = tf.add_paragraph()
        sp.space_before = Pt(12)
        sr = sp.add_run()
        sr.text = subtitle
        sr.font.size = Pt(18)
        sr.font.color.rgb = TEXT_GREY


def add_context_slide(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_navy_header_band(slide)
    _slide_title(slide, "Why this review")

    body = slide.shapes.add_textbox(Inches(0.7), Inches(1.4), Inches(12), Inches(5.5))
    tf = body.text_frame
    tf.word_wrap = True

    bullets = [
        "EAGLE's document generator currently produces eight Tier A acquisition documents as raw markdown — no template structure.",
        "Operators (contracting officers, card holders, SSAs) expect the standard federal layout when they print, sign, or attach these documents.",
        "Each draft template is ~50 lines: persona statement, required sections, small worked example, generation rules.",
        "All eight follow the uniform style of the existing 23 production prompts — review one, trust the rest.",
        "Approval here is the gate before the templates ship into doc_prompts.py and the document registry.",
    ]
    for i, b in enumerate(bullets):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after = Pt(10)
        rb = p.add_run()
        rb.text = "•  " + b
        rb.font.size = Pt(16)
        rb.font.color.rgb = TEXT_DARK


def add_review_criteria_slide(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_navy_header_band(slide)
    _slide_title(slide, "What we need from you")

    body = slide.shapes.add_textbox(Inches(0.7), Inches(1.4), Inches(12), Inches(5.5))
    tf = body.text_frame
    tf.word_wrap = True

    rows = [
        ("Authority", "Are the FAR / HHSAR / agency policy citations correct and complete?"),
        ("Sections", "Are all required sections present? Are any missing or unnecessary?"),
        ("Worked example", "Does the example reflect the actual format operators expect to see?"),
        ("Rules", "Are the generation guardrails appropriate? No missed regulatory requirements?"),
        ("Audience", "Is the right reviewer / signer named for this document?"),
    ]
    for i, (k, v) in enumerate(rows):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after = Pt(8)
        kr = p.add_run()
        kr.text = k + ":  "
        kr.font.size = Pt(16)
        kr.font.bold = True
        kr.font.color.rgb = NCI_NAVY
        vr = p.add_run()
        vr.text = v
        vr.font.size = Pt(16)
        vr.font.color.rgb = TEXT_DARK

    p = tf.add_paragraph()
    p.space_before = Pt(20)
    fr = p.add_run()
    fr.text = "Mark each: APPROVED · APPROVED WITH CHANGES · REJECTED — REWORK"
    fr.font.size = Pt(14)
    fr.font.italic = True
    fr.font.color.rgb = TEXT_GREY


def _slide_title(slide, title_text: str) -> None:
    tb = slide.shapes.add_textbox(Inches(0.7), Inches(0.6), Inches(12), Inches(0.7))
    tf = tb.text_frame
    p = tf.paragraphs[0]
    r = p.add_run()
    r.text = title_text
    r.font.size = Pt(26)
    r.font.bold = True
    r.font.color.rgb = NCI_NAVY


def add_template_slide(
    prs: Presentation,
    idx: int,
    slug: str,
    doc_name: str,
    authority: str,
    audience: str,
    sections: list[tuple[str, str]],
    worked_example: str,
) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_navy_header_band(slide)
    _slide_title(slide, f"Template {idx} of 8 — {doc_name}")

    # Slug + authority + audience (top-left meta block)
    meta = slide.shapes.add_textbox(Inches(0.7), Inches(1.35), Inches(12), Inches(1.0))
    mt = meta.text_frame
    mt.word_wrap = True
    for label, value in [("slug", slug), ("authority", authority), ("audience", audience)]:
        para = mt.add_paragraph() if mt.paragraphs[0].text else mt.paragraphs[0]
        para.space_after = Pt(2)
        rk = para.add_run()
        rk.text = f"{label}:  "
        rk.font.size = Pt(11)
        rk.font.bold = True
        rk.font.color.rgb = NCI_NAVY
        rv = para.add_run()
        rv.text = value
        rv.font.size = Pt(11)
        rv.font.color.rgb = TEXT_DARK

    # Sections list (left column)
    sec_box = slide.shapes.add_textbox(Inches(0.7), Inches(2.55), Inches(7.3), Inches(4.6))
    st = sec_box.text_frame
    st.word_wrap = True
    h = st.paragraphs[0]
    hr = h.add_run()
    hr.text = "Required Sections"
    hr.font.size = Pt(13)
    hr.font.bold = True
    hr.font.color.rgb = NCI_NAVY
    for n, (sec_name, sec_desc) in enumerate(sections, 1):
        p = st.add_paragraph()
        p.space_before = Pt(3)
        rn = p.add_run()
        rn.text = f"{n}. "
        rn.font.size = Pt(10)
        rn.font.bold = True
        rn.font.color.rgb = NCI_NAVY
        rname = p.add_run()
        rname.text = sec_name + " — "
        rname.font.size = Pt(10)
        rname.font.bold = True
        rname.font.color.rgb = TEXT_DARK
        rdesc = p.add_run()
        rdesc.text = sec_desc
        rdesc.font.size = Pt(10)
        rdesc.font.color.rgb = TEXT_GREY

    # Worked example (right column, light grey background)
    from pptx.enum.shapes import MSO_SHAPE
    bg = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(8.2), Inches(2.55), Inches(4.5), Inches(4.6))
    bg.fill.solid()
    bg.fill.fore_color.rgb = NCI_LIGHT
    bg.line.color.rgb = ACCENT_LINE
    bg.line.width = Pt(0.75)

    we = slide.shapes.add_textbox(Inches(8.4), Inches(2.7), Inches(4.2), Inches(4.3))
    wt = we.text_frame
    wt.word_wrap = True
    wh = wt.paragraphs[0]
    wr = wh.add_run()
    wr.text = "Worked Example"
    wr.font.size = Pt(13)
    wr.font.bold = True
    wr.font.color.rgb = NCI_NAVY
    wp = wt.add_paragraph()
    wp.space_before = Pt(5)
    wer = wp.add_run()
    wer.text = worked_example
    wer.font.size = Pt(10)
    wer.font.color.rgb = TEXT_DARK


def add_approval_matrix_slide(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_navy_header_band(slide)
    _slide_title(slide, "Approval Matrix")

    rows = [
        ("cor_designation", "COR Designation Letter"),
        ("sb_review", "HHS-653 Small Business Review"),
        ("section_889", "Section 889 Compliance"),
        ("section_508", "Section 508 Compliance"),
        ("priority_sources_checklist", "Required & Priority Sources Checklist"),
        ("subk_review", "Subcontracting Plan Review"),
        ("qasp", "Quality Assurance Surveillance Plan"),
        ("source_selection_plan", "Source Selection Plan"),
    ]
    headers = ["#", "Template", "APPROVED", "W/ CHANGES", "REJECTED", "Comments"]
    col_w = [Inches(0.4), Inches(4.5), Inches(1.4), Inches(1.4), Inches(1.4), Inches(3.0)]

    table_top = Inches(1.5)
    table_left = Inches(0.5)
    n_rows = len(rows) + 1
    table = slide.shapes.add_table(n_rows, len(headers), table_left, table_top,
                                    sum((c for c in col_w), Inches(0)), Inches(5.0)).table

    for ci, w in enumerate(col_w):
        table.columns[ci].width = w

    # Header row
    for ci, htext in enumerate(headers):
        cell = table.cell(0, ci)
        cell.text = ""
        p = cell.text_frame.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        r = p.add_run()
        r.text = htext
        r.font.size = Pt(11)
        r.font.bold = True
        r.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
        cell.fill.solid()
        cell.fill.fore_color.rgb = NCI_NAVY

    # Data rows
    for ri, (slug, name) in enumerate(rows, 1):
        for ci, val in enumerate([str(ri), name, "", "", "", ""]):
            cell = table.cell(ri, ci)
            cell.text = ""
            p = cell.text_frame.paragraphs[0]
            r = p.add_run()
            r.text = val
            r.font.size = Pt(10)
            r.font.color.rgb = TEXT_DARK
            if ci == 0 or ci >= 2:
                p.alignment = PP_ALIGN.CENTER
            if ri % 2 == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = NCI_LIGHT

    # Signature line below
    sig = slide.shapes.add_textbox(Inches(0.5), Inches(6.7), Inches(12), Inches(0.5))
    sf = sig.text_frame
    sp = sf.paragraphs[0]
    sr = sp.add_run()
    sr.text = "Reviewer signature: ___________________________   Name / Title: ___________________________   Date: __________"
    sr.font.size = Pt(11)
    sr.font.color.rgb = TEXT_GREY


def add_next_steps_slide(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_navy_header_band(slide)
    _slide_title(slide, "What happens after approval")

    body = slide.shapes.add_textbox(Inches(0.7), Inches(1.4), Inches(12), Inches(5.5))
    tf = body.text_frame
    tf.word_wrap = True

    steps = [
        "Approved templates → committed to server/app/doc_prompts.py as named constants.",
        "Each slug's DocSpec.system_prompt updated to reference the new constant in server/app/doc_registry.py.",
        "Module-load validator confirms registry consistency; unit tests run.",
        "If green, changes ship in the next deployment.",
        "APPROVED W/ CHANGES return to the platform team with the change list, then back here for re-review.",
        "REJECTED — REWORK gets a fresh draft and re-enters this review cycle.",
    ]
    for i, s in enumerate(steps):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after = Pt(10)
        rb = p.add_run()
        rb.text = f"{i+1}.  " + s
        rb.font.size = Pt(15)
        rb.font.color.rgb = TEXT_DARK


# ───────────────────────────────────────────────────────────────────────
# Per-template content (keep aligned with the markdown report)

TEMPLATES = [
    dict(
        slug="cor_designation",
        doc_name="COR Designation Letter",
        authority="FAR 1.602-2(d), HHSAR 301.604, OFPP 05-01",
        audience="CO (signs), COR nominee (acknowledges), supervisor (cc)",
        sections=[
            ("LETTERHEAD", "NCI/NIH letterhead, date, addressee block"),
            ("SUBJECT LINE", "Designation as COR — Contract [Number]"),
            ("DESIGNATION STATEMENT", "Appointment invoking FAR 1.602-2(d)"),
            ("CONTRACT IDENTIFICATION", "Number, contractor, PoP, value, type"),
            ("AUTHORIZED DUTIES", "Technical direction, inspection, CPARS"),
            ("LIMITATIONS", "Express prohibitions reserved to CO"),
            ("FAC-COR LEVEL", "Required level (I/II/III) for complexity"),
            ("REPORTING CADENCE", "Routine reports + immediate-notice events"),
            ("REVOCATION CLAUSE", "Revocable in writing; auto-terminates"),
            ("ACKNOWLEDGEMENT", "COR nominee accepts duties + limits"),
            ("SIGNATURE", "Contracting Officer + cc list"),
        ],
        worked_example=(
            "6. LIMITATIONS — Pursuant to FAR 1.602-2(d), the following are "
            "reserved exclusively to the CO and may NOT be undertaken by the COR "
            "for Contract HHSN272201800001I (Illumina BaseSpace maintenance):\n\n"
            "  1. Modify price, PoP, or scope of work.\n"
            "  2. Authorize work outside the SOW or any cost overrun.\n"
            "  3. Make any commitment that obligates the Government.\n"
            "  4. Release source-selection-sensitive information.\n"
            "  5. Settle claims or invoice disputes without CO approval.\n\n"
            "Any action exceeding this delegation is an unauthorized "
            "commitment and may create personal liability for the COR."
        ),
    ),
    dict(
        slug="sb_review",
        doc_name="HHS-653 Small Business Review",
        authority="FAR Part 19, HHSAR 319, FAR 19.502 (Rule of Two)",
        audience="Small Business Specialist, CO, SSA",
        sections=[
            ("FORM HEADER", "Estimated value, NAICS, size standard, activity"),
            ("REQUIREMENT SUMMARY", "What is being acquired + PoP"),
            ("SB PROGRAM ASSESSMENT", "8(a)/HUBZone/SDVOSB/WOSB/SB grid"),
            ("MARKET RESEARCH BASIS", "Sources sought, SAM searches, history"),
            ("NAICS / PSC ANALYSIS", "Confirm size standard + receipts/employee"),
            ("RECOMMENDATION", "Recommended set-aside or unrestricted"),
            ("SB SPECIALIST CONCURRENCE", "Name, signature, date"),
            ("CO DETERMINATION", "Name, signature, date"),
        ],
        worked_example=(
            "3. SB PROGRAM ASSESSMENT — NAICS 541512, $280K annual.\n\n"
            "  • 8(a) Suitable? NO — incumbent is large business; transition "
            "    not feasible within 60-day expiration.\n"
            "  • HUBZone? NO — sole-source product; no HUBZone vendor.\n"
            "  • SDVOSB / WOSB? NO — proprietary product, no qualified SB.\n"
            "  • Rule of Two (FAR 19.502)? NO — single-source product.\n\n"
            "Recommendation: Unrestricted sole source per FAR 6.302-1; "
            "rationale supported by sources sought — zero responses."
        ),
    ),
    dict(
        slug="section_889",
        doc_name="Section 889 Compliance Documentation",
        authority="FAR 4.21, FAR 52.204-24/25/26, NDAA FY2019 § 889",
        audience="CO, vendor (representations), agency leadership (waivers)",
        sections=[
            ("MEMO HEADER", "Requirement, contract #, value, CO, date"),
            ("APPLICABILITY", "Part A (use) and/or Part B (essential component)"),
            ("COVERED ENTITIES", "Huawei, ZTE, Hytera, Hikvision, Dahua + affiliates"),
            ("REPRESENTATION REVIEW", "FAR 52.204-24 + 52.204-26 + date checked"),
            ("SAM.gov VERIFICATION", "Lookup date + result + screenshot"),
            ("WAIVER STATUS", "FAR 4.2104 cite + expiration if applicable"),
            ("REPORTING OBLIGATIONS", "52.204-25(d) 1-day / 10-day"),
            ("FINDING", "Award is/is not permitted under FAR 4.21"),
            ("SIGNATURE", "Contracting officer, date"),
        ],
        worked_example=(
            "4. REPRESENTATION REVIEW — Illumina Inc., UEI ABC123\n\n"
            "  • FAR 52.204-24 (offer-level rep): Submitted 2026-05-05. "
            "    Section A (covered telecom): NO. Section B (substantial "
            "    or essential component): NO.\n"
            "  • FAR 52.204-26 (annual rep): Verified active in SAM.gov "
            "    on 2026-05-05; SAM annual rep dated 2026-02-14.\n"
            "  • Discrepancies: None.\n"
            "  • Waiver: Not applicable.\n\n"
            "Determination: Section 889 representations support award."
        ),
    ),
    dict(
        slug="section_508",
        doc_name="Section 508 Compliance Statement",
        authority="29 USC 794d, 36 CFR 1194 (Revised 508), FAR 39.2",
        audience="508 Coordinator, CO, vendor (VPAT submission)",
        sections=[
            ("APPLICABILITY", "Does this acquisition include EIT?"),
            ("PRODUCT TYPE CHECKLIST", "Software/Web/Telecom/Video/SC/Desktop/E-Doc"),
            ("EXCEPTION DETERMINATION", "Nat-sec / incidental / micro / fund-alt / undue"),
            ("VPAT/ACR STATUS", "Vendor accessibility conformance report"),
            ("CONTRACT LANGUAGE", "SOW + eval criteria + clause checklist"),
            ("508 COORDINATOR REVIEW", "Reviewer name, date, finding"),
            ("REQUIRED CLAUSES", "FAR 52.239-1, HHSAR 352.239-73"),
            ("SIGNATURE", "Contracting officer, 508 coordinator"),
        ],
        worked_example=(
            "2. PRODUCT TYPE CHECKLIST — Cloud genomics analysis platform\n\n"
            "  ☑ Software — Yes (SaaS web app)\n"
            "  ☑ Web — Yes (browser-based UI)\n"
            "  ☐ Telecom — No\n"
            "  ☐ Video — No\n"
            "  ☐ Self-Contained / Closed Products — No\n"
            "  ☐ Desktop / Mobile Devices — No\n"
            "  ☑ Electronic Documents — Yes (PDF reports, exports)\n\n"
            "Pre-checked from contract context (description = 'cloud-based "
            "genomics analysis platform with reporting'). Vendor must submit "
            "current VPAT-ACR for Software, Web, and Electronic Documents."
        ),
    ),
    dict(
        slug="priority_sources_checklist",
        doc_name="Required & Priority Sources Checklist",
        authority="FAR Part 8, HHS supplemental, FAR 8.002(a)(1)",
        audience="Card holder / CO; AbilityOne / UNICOR; Strategic Sourcing",
        sections=[
            ("HEADER", "Title, value, requestor, CO, date"),
            ("REQUIREMENT SUMMARY", "2-3 sentence description"),
            ("PRIORITY-OF-USE TABLE", "8 sources in FAR Part 8 order"),
            ("AbilityOne / UNICOR DETERMINATION", "Mandatory if on Procurement List"),
            ("STRATEGIC SOURCING / BIC", "BIC vehicle note + rationale if not used"),
            ("FINDING", "Selected source + cite to FAR 8.xxx subsection"),
            ("SIGNATURE", "Card holder / contracting officer"),
        ],
        worked_example=(
            "3. PRIORITY-OF-USE TABLE — Office supplies, $7,500 micro-purchase\n\n"
            "  1. Agency inventory — Checked. None available in storeroom.\n"
            "  2. Excess from other agencies — N/A (commodity, not surplus).\n"
            "  3. UNICOR — Checked AbilityOne PL ✓. Item NOT on PL.\n"
            "  4. AbilityOne — Checked AbilityOne PL ✓. Item NOT on PL.\n"
            "  5. Wholesale supply — N/A (no DLA / GSA Global Supply match).\n"
            "  6. Mandatory FSS — Checked GSA Schedule 36 ✓. Available.\n\n"
            "Selection: Mandatory FSS (FAR 8.002(a)(1)(vi))."
        ),
    ),
    dict(
        slug="subk_review",
        doc_name="Subcontracting Plan Review",
        authority="FAR 19.705-4, FAR 52.219-9, HHSAR 352.219",
        audience="CO, SBS, SBA PCR (when applicable), SSA",
        sections=[
            ("MEMO HEADER", "Date, contractor, contract #, plan version, reviewer"),
            ("PLAN TYPE", "Individual / Master / Commercial"),
            ("ADEQUACY CHECKLIST", "10-row table — each FAR 52.219-9 element"),
            ("GOAL REASONABLENESS", "Goals vs prior-year + industry benchmark"),
            ("FINDINGS", "Numbered list of deficiencies / changes required"),
            ("SBA PCR REVIEW", "PCR concurrence when applicable per 19.705-4(d)"),
            ("RECOMMENDATION", "Approve / Approve w/ revisions / Return / Reject"),
            ("SIGNATURE", "Contracting officer, date"),
        ],
        worked_example=(
            "3. ADEQUACY CHECKLIST — Excerpt (3 of 10 rows)\n\n"
            "  | # | Element                   | Present | Adequate | Notes |\n"
            "  | 1 | Goals (SB/SDB/WOSB/HZ/...) |   ☑    |    ☑    | 28% SB|\n"
            "  | 2 | Products/svc subcontracted |   ☑    |    ☑    | NAICS |\n"
            "  | 3 | Method for developing goals|   ☑    |    ☐    | thin  |\n\n"
            "Element 3 deficient: rationale references 'industry standard' "
            "without citing data. Return for revision with prior-year "
            "achievement + SBA market research evidence."
        ),
    ),
    dict(
        slug="qasp",
        doc_name="Quality Assurance Surveillance Plan",
        authority="FAR 46.401, FAR 37.6 PBA, NIH PBA handbook",
        audience="COR (executes), CO (approves + remedies), contractor",
        sections=[
            ("PURPOSE", "Methods + resources for performance monitoring"),
            ("SCOPE", "Linked PWS + objectives covered + period"),
            ("ROLES", "COR / CO / Tech POC / alternate / surveillance team"),
            ("OBJECTIVES TABLE", "1 row per PWS objective: AQL + method"),
            ("METHODS CATALOG", "100% / random / periodic / complaint / audit"),
            ("AQL DEFINITIONS", "Measurable units + calculation method"),
            ("INCENTIVES", "Award fee / cure / show-cause; omit if FFP"),
            ("DOCUMENTATION", "Logs, monthly assessments, CPARS"),
            ("SIGNATURE", "COR + Contracting Officer"),
        ],
        worked_example=(
            "4. OBJECTIVES TABLE — Genomic data analysis platform\n\n"
            "  | # | Objective         | PWS  | AQL              | Method  | Freq |\n"
            "  | 1 | Platform uptime   | §5.1 | ≥99.5% monthly   | 100% telemetry | mo |\n"
            "  | 2 | Tier-2 ticket SLA | §5.3 | ≥95% in 4 hrs    | random sample 20/mo | mo |\n"
            "  | 3 | Compliance report | §5.7 | by 5th biz day; 0 critical findings | 100% on receipt | mo |\n\n"
            "Method assignment: 1+3 use 100% (telemetry / single deliverable), "
            "2 uses random sampling (>500 tickets/mo makes 100% impractical)."
        ),
    ),
    dict(
        slug="source_selection_plan",
        doc_name="Source Selection Plan (SSP)",
        authority="FAR 15.3, HHSAR 315.3, FAR 15.300-308",
        audience="SSA, CO, TET, Cost/Price, Legal counsel",
        sections=[
            ("ACQUISITION OVERVIEW", "Title, value, PoP, NAICS, set-aside"),
            ("SSA DESIGNATION", "Source Selection Authority + decision authority"),
            ("SOURCE SELECTION TEAM", "SSA / CO / TET / Cost / Legal / SBS"),
            ("EVALUATION FACTORS", "Descending importance + subfactors"),
            ("METHODOLOGY", "Color/adjectival rating + risk + PP confidence"),
            ("PRICE/COST APPROACH", "FAR 15.404 techniques to be used"),
            ("COMMUNICATIONS", "Clarifications vs discussions per 15.306"),
            ("AWARD FRAMEWORK", "Best-value tradeoff vs LPTA"),
            ("DOCUMENTATION", "SSDD, debriefings, source selection statement"),
            ("SCHEDULE", "Solicitation through award milestones"),
            ("APPROVALS", "SSA, CO, TET Chair signature block"),
        ],
        worked_example=(
            "4. EVALUATION FACTORS — Cloud genomics platform\n\n"
            "  Factor 1 — Technical / Mission Suitability (most important)\n"
            "    1a. Solution architecture\n"
            "    1b. Security & compliance (FedRAMP, FISMA)\n"
            "    1c. Implementation approach\n"
            "  Factor 2 — Past Performance (significantly less than Factor 1)\n"
            "  Factor 3 — Price / Cost (lowest)\n\n"
            "Relative-importance statement: Technical and Past Performance "
            "combined are significantly more important than Price."
        ),
    ),
]

# ───────────────────────────────────────────────────────────────────────

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

add_title_slide(prs)
add_context_slide(prs)
add_review_criteria_slide(prs)
add_section_header_slide(prs, "Eight Templates", "One slide per template — review structure, sections, and example")

for i, t in enumerate(TEMPLATES, 1):
    add_template_slide(
        prs, idx=i,
        slug=t["slug"], doc_name=t["doc_name"], authority=t["authority"],
        audience=t["audience"], sections=t["sections"], worked_example=t["worked_example"],
    )

add_approval_matrix_slide(prs)
add_next_steps_slide(prs)

prs.save(str(OUT_PPTX))
print(f"wrote {OUT_PPTX}")
print(f"size: {OUT_PPTX.stat().st_size:,} bytes")
print(f"slides: {len(prs.slides)}")
